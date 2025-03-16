from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from json_data import fetch_exam_data_default
import requests
import os
import re
import random
import string

def fetch_exam_data(exam_id):
    url = f"http://127.0.0.1:8000/api/v3/web/exam-question-answer-json/{exam_id}"
    response = requests.get(url)
    if response.status_code == 200:
        return response.json()
    else:
        raise Exception(f"Failed to fetch exam data: {response.status_code}")

# Function to generate PowerPoint
def generate_question_answer_ppt(exam):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    serial_no = 1

    for question in exam["questions"]:
        slide = prs.slides.add_slide(slide_layout)
        title_content = f"{serial_no}. {clean_html_entities(question['title'])}"
        img_link = extract_image_url(title_content)
        offset_y = 1
        if img_link:
            try:
                img_path = download_image(img_link)
                discussion_slide.shapes.add_picture(img_path, Inches(1), Inches(2), height=Inches(3.5))
                offset_y = 4
            except Exception:
                pass
        if slide.shapes.title:
            sp = slide.shapes.title
            slide.shapes._spTree.remove(sp._element)

        text_frame = slide.placeholders[1].text_frame
        text_frame.text = title_content
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.size = Pt(22)
        text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        for option in question["question_answers"]:
            p = text_frame.add_paragraph()
            p.text = clean_html_entities(option["answer"])
            p.font.size = Pt(18)
            p.alignment = PP_ALIGN.LEFT

        if len(question["reference_books"]):
            ref_book = ""
            # print(question["reference_books"])
            for ref in question["reference_books"]:
                if "reference_book_id" in ref and "page_no" in ref:
                    ref_book += f"[Ref: {ref['reference_book']['name']}/P-{ref['page_no']}] "
        
        # Create a new shape (textbox) for the answer
        left = Inches(1)
        top = Inches(1)
        width = Inches(6)
        height = Inches(1)

        # This is now a separate shape for the answer
        answer_shape = slide.shapes.add_textbox(left, top, width, height)
        text_frame = answer_shape.text_frame
        ans_p = text_frame.add_paragraph()
        ans_p.text = "Ans: This is a separate textbox"
        ans_p.font.bold = True
        ans_p.font.size = Pt(18)

        # Now we can apply animation to the whole shape!
        fade_xml = f"""
        <p:anim xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
            <p:cBhvr>
                <p:cTn id="1" dur="500" fill="hold">
                    <p:stCondLst>
                        <p:cond delay="0"/>
                    </p:stCondLst>
                </p:cTn>
                <p:tgtEl>
                    <p:spTgt spid="{answer_shape.shape_id}"/>
                </p:tgtEl>
            </p:cBhvr>
            <p:animEffect transition="in" filter="fade"/>
        </p:anim>
        """

        shape_element = answer_shape.element
        shape_element.getparent().append(parse_xml(fade_xml))
        
        serial_no += 1
        if (question.get("discussion", "") and question.get("discussion", "").strip()):
            discussion_slide = prs.slides.add_slide(slide_layout)
            # Remove the title shape if it exists
            if discussion_slide.shapes.title:
                sp = discussion_slide.shapes.title
                discussion_slide.shapes._spTree.remove(sp._element)
                
            discussion_frame = discussion_slide.placeholders[1].text_frame

            if ref_book.strip() == '':
                ref_content = question.get('reference', '') if question.get('reference', '') else ''
                ref_content = clean_html_entities(ref_content)
                ans_p = discussion_frame.add_paragraph()
                ans_p.text = ref_content
                ans_p.font.bold = True
                ans_p.font.size = Pt(18)
                ans_p.alignment = PP_ALIGN.LEFT
           
        
        if question.get("discussion", "") and question.get("discussion", "").strip():
            img_link = extract_image_url(question["discussion"])
            offset_y = 1

            if img_link:
                try:
                    img_path = download_image(img_link)
                    discussion_slide.shapes.add_picture(img_path, Inches(1), Inches(2), height=Inches(3.5))
                    offset_y = 4
                except Exception:
                    pass

            discussion_frame.top = Inches(1)
            discussion_p = discussion_frame.add_paragraph()
            discussion_p.text = clean_html_entities(question["discussion"])
            discussion_p.font.size = Pt(18)
            discussion_p.alignment = PP_ALIGN.LEFT

    exam_name = re.sub(r"[^a-zA-Z0-9]", "_", exam["name"] or "presentation")
    ppt_file_name = f"{exam_name}.pptx"
    prs.save(ppt_file_name)
    return ppt_file_name

def clean_html_entities(text):
    text = re.sub(r'&.*?;', "", text)
    text = re.sub(r'\r', "", text)
    return re.sub(r"<.*?>", "", text)

def extract_image_url(html):
    match = re.search(r'<img[^>]+src=["\']?([^"\'>]+)["\']?', html, re.I)
    return match.group(1) if match else None

def download_image(url):
    response = requests.get(url, stream=True)
    if response.status_code == 200:
        file_name = f"{''.join(random.choices(string.ascii_letters, k=10))}.jpg"
        file_path = os.path.join(os.getcwd(), file_name)
        with open(file_path, "wb") as file:
            for chunk in response.iter_content(1024):
                file.write(chunk)
        return file_path
    else:
        raise Exception(f"Failed to download image: {url}")

if __name__ == "__main__":
    exam_id = input("Enter Exam ID: ")
    try:
        # exam_data = fetch_exam_data(exam_id)
        exam_data = fetch_exam_data_default()
        ppt_file = generate_question_answer_ppt(exam_data)
        print(f"PowerPoint generated: {ppt_file}")
    except Exception as e:
        print(f"Error: {e}")

