import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import requests
import re
import threading
import os
import re
import random
import string
# Function to fetch exam data
def fetch_exam_data(exam_id):
    url = f"https://admin.genesisedu.info/api/v3/web/exam-question-answer-json/{exam_id}"
    response = requests.get(url)
    if response.status_code == 200:
        return response.json()
    else:
        raise Exception(f"Failed to fetch exam data: {response.status_code}")

# Function to generate PowerPoint
def generate_question_answer_ppt(exam):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    serial_no = 1

    for question in exam["questions"]:
        slide = prs.slides.add_slide(slide_layout)
        title_content = f"{serial_no}. {clean_html_entities(question['title'])}"

        if slide.shapes.title:
            sp = slide.shapes.title
            slide.shapes._spTree.remove(sp._element)

        text_frame = slide.placeholders[1].text_frame
        text_frame.text = title_content
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.size = Pt(22)
        text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

        for option in question["question_answers"]:
            p = text_frame.add_paragraph()
            p.text = clean_html_entities(option["answer"])
            p.font.size = Pt(18)
            p.alignment = PP_ALIGN.LEFT

        ans_p = text_frame.add_paragraph()
        ans_p.text = f"Ans: {question.get('answer_script', '')}"
        ans_p.font.bold = True
        ans_p.font.size = Pt(18)
        ans_p.alignment = PP_ALIGN.LEFT

        serial_no += 1
        if (question.get("discussion", "") and question.get("discussion", "").strip()) or len(question["reference_books"]):
            discussion_slide = prs.slides.add_slide(slide_layout)
            # Remove the title shape if it exists
            if discussion_slide.shapes.title:
                sp = discussion_slide.shapes.title
                discussion_slide.shapes._spTree.remove(sp._element)
                
            discussion_frame = discussion_slide.placeholders[1].text_frame
        # Process reference books
        if len(question["reference_books"]):
            ref_book = ""
            # print(question["reference_books"])
            for ref in question["reference_books"]:
                if "reference_book_id" in ref and "page_no" in ref:
                    ref_book += f"[Ref: {ref['reference_book']['name']}/P-{ref['page_no']}] "

            if ref_book :
                ref_content = ref_book
            else:
                ref_content = question.get('reference', '') if question.get('reference', '') else ''
                
            ref_content = clean_html_entities(ref_content)
            ans_p = discussion_frame.add_paragraph()
            ans_p.text = ref_content
            ans_p.font.bold = True
            ans_p.font.size = Pt(18)
            ans_p.alignment = PP_ALIGN.LEFT
           
        
        if question.get("discussion", "") and question.get("discussion", "").strip():
            img_link = extract_image_url(question["discussion"])
            offset_y = 1

            if img_link:
                try:
                    img_path = download_image(img_link)
                    discussion_slide.shapes.add_picture(img_path, Inches(1), Inches(2), height=Inches(3.5))
                    offset_y = 4
                except Exception:
                    pass

            discussion_frame.top = Inches(1)
            discussion_p = discussion_frame.add_paragraph()
            discussion_p.text = clean_html_entities(question["discussion"])
            discussion_p.font.size = Pt(18)
            discussion_p.alignment = PP_ALIGN.LEFT
            
    exam_name = re.sub(r"[^a-zA-Z0-9]", "_", exam["name"] or "presentation")
    ppt_file_name = filedialog.asksaveasfilename(defaultextension=".pptx",
                                                 filetypes=[("PowerPoint Files", "*.pptx")],
                                                 initialfile=f"{exam_name}.pptx")
    if ppt_file_name:
        prs.save(ppt_file_name)
        messagebox.showinfo("Success", f"PowerPoint saved as: {ppt_file_name}")

# Utility Functions
def clean_html_entities(text):
    text = re.sub(r'&.*?;', "", text)
    text = re.sub(r'\r', "", text)
    return re.sub(r"<.*?>", "", text)

def extract_image_url(html):
    match = re.search(r'<img[^>]+src=["\']?([^"\'>]+)["\']?', html, re.I)
    return match.group(1) if match else None

def download_image(url):
    response = requests.get(url, stream=True)
    if response.status_code == 200:
        file_name = f"{''.join(random.choices(string.ascii_letters, k=10))}.jpg"
        file_path = os.path.join(os.getcwd(), file_name)
        with open(file_path, "wb") as file:
            for chunk in response.iter_content(1024):
                file.write(chunk)
        return file_path
    else:
        raise Exception(f"Failed to download image: {url}")
    
# Function to handle button click
def start_generation():
    exam_id = exam_id_entry.get().strip()
    if not exam_id:
        messagebox.showerror("Error", "Please enter an Exam ID")
        return

    def worker():
        try:
            exam_id_entry.delete(0, tk.END)
            exam_data = fetch_exam_data(exam_id)
            generate_question_answer_ppt(exam_data)
        except Exception as e:
            messagebox.showerror("Error", str(e))

    # Run in a separate thread to keep the GUI responsive
    thread = threading.Thread(target=worker)
    thread.start()

# GUI Setup
root = tk.Tk()
root.title("PowerPoint Generator")
root.geometry("400x230")
root.resizable(False, False)  # Fixed window size
root.configure(bg="lightblue")  # Light background color

# Title Label
title_label = tk.Label(root, text="Exam ID", font=("Arial", 20, "bold"), fg="indigo",  bg="lightblue")
title_label.pack(pady=10)

# Entry Field
exam_id_entry = tk.Entry(root, width=20, font=("Arial", 15), bd=3)
exam_id_entry.pack(padx=2, pady=4)

# Generate Button
generate_button = tk.Button(
    root, 
    text="Generate PowerPoint", 
    font=("Arial", 12),
    fg="white",
    bg="indigo",  # Corrected color
    padx=8, 
    pady=6,
    cursor="hand2",
    command=start_generation
)
generate_button.pack(pady=16)

root.mainloop()